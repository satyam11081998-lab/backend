"""
Local Deck Content Extractor & Preprocessor.

Extracts text, slide titles, tables, page counts, and numerical tokens locally from PDF and PPTX files
without sending raw files to external LLMs. Also performs deterministic signal mining on file names
and folder hierarchies.
"""

import hashlib
import os
import re
from typing import Any, Dict, List, Optional, Set, Tuple

import pypdfium2 as pdfium

try:
    import pptx
except ImportError:
    pptx = None

from services.deck_taxonomy import (
    BSCHOOLS,
    CASE_TYPES,
    KNOWN_COMPETITIONS,
    RESULTS,
    ROUND_TYPES,
)


def compute_file_hash(file_path_or_bytes: Any) -> str:
    """Compute deterministic SHA-256 hash of a file or byte buffer."""
    hasher = hashlib.sha256()
    if isinstance(file_path_or_bytes, (bytes, bytearray)):
        hasher.update(file_path_or_bytes)
    else:
        with open(file_path_or_bytes, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                hasher.update(chunk)
    return hasher.hexdigest()


_THOUSANDS_SEP = re.compile(r"(?<=\d),(?=\d{3}(?:\D|$))")
_WHITESPACE_RE = re.compile(r"[ \t]+")
_MULTI_NEWLINE_RE = re.compile(r"\n{3,}")


def extract_numbers_from_text(text: str) -> Set[str]:
    """Extract normalized numbers from text for zero-hallucination verification."""
    cleaned = _THOUSANDS_SEP.sub("", text or "")
    out: Set[str] = set()
    for token in re.findall(r"\d+(?:\.\d+)?", cleaned):
        try:
            val = float(token)
        except ValueError:
            out.add(token)
            continue
        out.add(str(int(val)) if val == int(val) else repr(val).rstrip("0").rstrip("."))
    return out


def normalize_extracted_text(text: str) -> str:
    """Normalize whitespaces, tabs, and unprintable characters in extracted text."""
    if not text:
        return ""
    # Replace unicode non-standard dashes and quotes
    text = text.replace("’", "'").replace("‘", "'").replace("“", '"').replace("”", '"')
    text = text.replace("—", "-").replace("–", "-")
    lines = [_WHITESPACE_RE.sub(" ", line).strip() for line in text.splitlines()]
    # Remove empty line clumps
    normalized = "\n".join(lines)
    return _MULTI_NEWLINE_RE.sub("\n\n", normalized).strip()


def extract_pdf_deck(file_path: str) -> Dict[str, Any]:
    """Extract slides, text, and structure from a PDF file using pypdfium2."""
    with open(file_path, "rb") as f:
        pdf_bytes = f.read()

    pdf = pdfium.PdfDocument(pdf_bytes)
    slide_count = len(pdf)
    slides: List[Dict[str, Any]] = []
    full_text_parts: List[str] = []

    for i in range(slide_count):
        page_num = i + 1
        page = pdf[i]
        textpage = None
        raw_text = ""
        try:
            textpage = page.get_textpage()
            raw_text = textpage.get_text_range() or ""
        finally:
            if textpage is not None:
                try:
                    textpage.close()
                except Exception:
                    pass
            try:
                page.close()
            except Exception:
                pass

        clean_text = normalize_extracted_text(raw_text)
        lines = [line.strip() for line in clean_text.splitlines() if line.strip()]
        title = lines[0] if lines else f"Slide {page_num}"
        # If the first line is very short or generic like 'Team', check line 2
        if len(lines) > 1 and len(title) < 6 and lines[1]:
            title = f"{title}: {lines[1]}"

        slides.append({
            "slide_number": page_num,
            "title": title[:120],
            "text": clean_text,
        })
        if clean_text:
            full_text_parts.append(f"--- Slide {page_num}: {title} ---\n{clean_text}")

    try:
        pdf.close()
    except Exception:
        pass

    full_text = "\n\n".join(full_text_parts)
    return {
        "slide_count": slide_count,
        "slides": slides,
        "full_text": full_text,
        "numbers": extract_numbers_from_text(full_text),
    }


def extract_pptx_deck(file_path: str) -> Dict[str, Any]:
    """Extract slides, text, shapes, and notes from PPTX using python-pptx."""
    if pptx is None:
        raise RuntimeError("python-pptx library is not installed.")

    prs = pptx.Presentation(file_path)
    slides: List[Dict[str, Any]] = []
    full_text_parts: List[str] = []

    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        slide_texts: List[str] = []
        slide_title = ""

        # Check if presentation slide has standard title
        try:
            if slide.shapes.title and slide.shapes.title.text:
                slide_title = slide.shapes.title.text.strip()
        except Exception:
            pass

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        slide_texts.append(t)
            elif shape.has_table:
                # Extract tabular data as structured text
                for row in shape.table.rows:
                    row_cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if row_cells:
                        slide_texts.append(" | ".join(row_cells))

        # Check speaker notes if present
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Speaker Notes]: {notes}")
        except Exception:
            pass

        raw_text = "\n".join(slide_texts)
        clean_text = normalize_extracted_text(raw_text)
        lines = [line.strip() for line in clean_text.splitlines() if line.strip()]

        if not slide_title:
            slide_title = lines[0] if lines else f"Slide {page_num}"

        slides.append({
            "slide_number": page_num,
            "title": slide_title[:120],
            "text": clean_text,
        })
        if clean_text:
            full_text_parts.append(f"--- Slide {page_num}: {slide_title} ---\n{clean_text}")

    full_text = "\n\n".join(full_text_parts)
    return {
        "slide_count": len(prs.slides),
        "slides": slides,
        "full_text": full_text,
        "numbers": extract_numbers_from_text(full_text),
    }


def mine_path_and_filename_signals(file_path: str) -> Dict[str, Any]:
    """
    Extract high-value deterministic clues from folder hierarchy and file names.
    e.g. 'Accenture B-School Challenge- National Finalist/Aviators_Pranjal Tyagi.pdf'
    """
    norm_path = os.path.normpath(file_path).replace("\\", "/")
    path_parts = [p.strip() for p in norm_path.split("/") if p.strip()]
    filename = os.path.basename(file_path)
    base_name, _ = os.path.splitext(filename)

    combined_path_text = " ".join(path_parts).lower()

    signals: Dict[str, Any] = {
        "raw_path": norm_path,
        "filename": filename,
        "path_parts": path_parts,
        "detected_competition": None,
        "detected_company": None,
        "detected_result": None,
        "detected_round_type": None,
        "detected_case_type": None,
        "detected_year": None,
        "detected_college": None,
        "detected_team": None,
        "source_kind": "corporate",
    }

    # 1. Year detection (e.g. 2020..2030 or '24, '25)
    year_match = re.search(r"\b(202[0-9]|201[5-9])\b", combined_path_text)
    if year_match:
        signals["detected_year"] = int(year_match.group(1))
    else:
        # short year like '24 or `25
        short_year = re.search(r"['`’](2[0-9])\b", combined_path_text)
        if short_year:
            signals["detected_year"] = 2000 + int(short_year.group(1))

    # 2. Result detection
    for r in RESULTS:
        r_clean = r.lower()
        if r_clean in combined_path_text or r_clean.replace(" ", "") in combined_path_text.replace(" ", ""):
            signals["detected_result"] = r
            break

    # Also check variations like 'semi finalist', 'runner up', 'winners'
    if not signals["detected_result"]:
        if "national winner" in combined_path_text or "winner" in combined_path_text:
            signals["detected_result"] = "National Winner"
        elif "1st runner" in combined_path_text or "first runner" in combined_path_text:
            signals["detected_result"] = "National 1st Runner Up"
        elif "2nd runner" in combined_path_text or "second runner" in combined_path_text:
            signals["detected_result"] = "National 2nd Runner Up"
        elif "runner up" in combined_path_text or "runner-up" in combined_path_text:
            signals["detected_result"] = "National 1st Runner Up"
        elif "semi finalist" in combined_path_text or "semi-finalist" in combined_path_text or "semis" in combined_path_text:
            signals["detected_result"] = "National Semi Finalist"
        elif "finalist" in combined_path_text or "national finalist" in combined_path_text:
            signals["detected_result"] = "National Finalist"
        elif "shortlisted" in combined_path_text:
            signals["detected_result"] = "Shortlisted"
        elif "campus winner" in combined_path_text:
            signals["detected_result"] = "Campus Winner"

    # 3. Competition & Company detection
    for key, info in KNOWN_COMPETITIONS.items():
        if key in combined_path_text:
            signals["detected_competition"] = info["canonical_name"]
            signals["detected_company"] = info["company"]
            signals["source_kind"] = info["source_kind"]
            if not signals.get("detected_case_type"):
                signals["detected_case_type"] = info.get("default_case_type")
            break

    # 4. College / B-School detection
    for key, college_name in BSCHOOLS.items():
        key_pattern = re.escape(key)
        # Check standard word boundary OR concatenated in tokens e.g. MDIGurgaon or IMI_NewDelhi
        if re.search(r"\b" + key_pattern + r"\b", combined_path_text) or key.replace(" ", "") in combined_path_text.replace("_", "").replace("-", ""):
            signals["detected_college"] = college_name
            break

    # 5. Round type detection
    for rnd in ROUND_TYPES:
        if rnd in combined_path_text:
            signals["detected_round_type"] = rnd
            break
    if not signals["detected_round_type"]:
        if "semi" in combined_path_text:
            signals["detected_round_type"] = "semi-final"
        elif "final" in combined_path_text or "grand finale" in combined_path_text:
            signals["detected_round_type"] = "finale"
        elif "campus" in combined_path_text:
            signals["detected_round_type"] = "campus"
        elif "screening" in combined_path_text or "submission" in combined_path_text:
            signals["detected_round_type"] = "screening"

    # 6. Team name detection from filename
    team_match = re.search(r"(?:team\s+|team_)([a-zA-Z0-9_\s]+?)(?:[-_.]|$)", base_name, re.IGNORECASE)
    if team_match:
        signals["detected_team"] = team_match.group(1).replace("_", " ").strip()

    return signals


def extract_deck(file_path: str) -> Dict[str, Any]:
    """
    Main entry point for local deck extraction.
    Returns complete local structural, textual, numerical, and path metadata.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_size = os.path.getsize(file_path)
    file_hash = compute_file_hash(file_path)
    filename = os.path.basename(file_path)
    ext = os.path.splitext(filename)[1].lstrip(".").lower()

    if ext not in ("pdf", "pptx", "ppt", "xlsx"):
        raise ValueError(f"Unsupported deck format: .{ext}")

    if ext == "pdf":
        extracted = extract_pdf_deck(file_path)
    elif ext == "pptx":
        extracted = extract_pptx_deck(file_path)
    else:
        # Fallback for ppt or others
        extracted = {
            "slide_count": 0,
            "slides": [],
            "full_text": "",
            "numbers": set(),
        }

    path_signals = mine_path_and_filename_signals(file_path)

    return {
        "file_hash": file_hash,
        "file_path": os.path.abspath(file_path),
        "original_filename": filename,
        "file_size": file_size,
        "file_type": ext,
        "slide_count": extracted["slide_count"],
        "slides": extracted["slides"],
        "full_text": extracted["full_text"],
        "numbers": extracted["numbers"],
        "path_signals": path_signals,
    }
